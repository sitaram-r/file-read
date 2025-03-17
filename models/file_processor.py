import magic
import mammoth
import olefile
import subprocess
import glob
import tempfile
import os
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation

def process_file(file):
    file_info = {
        "file_name": file.filename  # ✅ Get exact file name
    }

    # Detect file_type using magic
    mime = magic.Magic(mime=True)
    mime_type = mime.from_buffer(file.read(2048))
    file.seek(0)  # Reset pointer

    file_info["mime_type"] = mime_type
    # print("Detected MIME Type:", mime_type)
    # Process PDF
    if mime_type == "application/pdf":
        reader = PdfReader(file)
        file_info["file_type"] = "pdf"
        file_info["page_count"] = len(reader.pages)
        file_info["text_preview"] = reader.pages[0].extract_text()[:3000] if reader.pages else "No text found"

    # Process DOCX
    elif mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        doc = Document(file)
        file_info["file_type"] = "docx"
        file_info["paragraph_count"] = len(doc.paragraphs)
        file_info["text_preview"] = " ".join([p.text for p in doc.paragraphs[:5]]) 

    # Process DOC (Older Word Files)
    elif mime_type == "application/msword":
        try:
            # Use `mammoth` to extract text from .doc
            result = mammoth.extract_raw_text(file)
            file_info["text_preview"] = result.value[:5000] if result else "No text found"
        except:
            try:
                # Use `olefile` as a backup for .doc
                if olefile.isOleFile(file):
                    ole = olefile.OleFileIO(file)
                    streams = ole.listdir()
                    file_info["text_preview"] = "Contains structured streams: " + str(streams)
                else:
                    file_info["Error"] = "Older .doc files are not fully supported. Convert to .docx."
            except:
                file_info["Error"] = "Failed to read .doc file."

    # Process PPTX
    elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        ppt = Presentation(file)
        file_info["slide_count"] = len(ppt.slides)
        file_info["file_type"] = "pptx"
        file_info["text_preview"] = " ".join(
            [shape.text for slide in ppt.slides for shape in slide.shapes if hasattr(shape, "text")][:5000]
        )

   # ✅ Read old PPT using catppt
    elif mime_type == "application/vnd.ms-powerpoint":
        try:
            result = subprocess.run(["catppt", "/path/to/tempfile.ppt"], stdout=subprocess.PIPE, text=True)
            file_info["text_preview"] = result.stdout[:5000] if result.stdout else "No text found"
        except:
            file_info["Error"] = "Failed to read .ppt file. Convert to .pptx."

    # ✅ Handle OLE Storage Format (.doc or .ppt)
    if mime_type == "application/x-ole-storage":
        if olefile.isOleFile(file):
            file.seek(0)

            # 🔹 If it's a .ppt file (Save temporarily & read with `soffice`)
            try:
                with tempfile.NamedTemporaryFile(delete=False, suffix=".ppt") as temp_file:
                    temp_file.write(file.read())
                    temp_file_path = temp_file.name

                soffice_path = "/opt/homebrew/bin/soffice"

                # 🔹 Convert PPT to PDF
                subprocess.run([
                    soffice_path, "--headless", "--convert-to", "pdf", "--outdir", "/tmp", temp_file_path
                ], stdout=subprocess.PIPE, stderr=subprocess.PIPE)

                # 🔹 Find the actual converted PDF (since LibreOffice renames it)
                pdf_files = glob.glob("/tmp/*.pdf")
                if not pdf_files:
                    file_info["Error"] = "Failed to find the converted PDF in /tmp"
                    return file_info

                # 🔹 Pick the latest converted file (should match our PPT)
                pdf_file_path = max(pdf_files, key=os.path.getctime)
                file_info["file_type"] = "ppt"
                # 🔹 Check slide count
                file_info["slide_count"] = count_pdf_pages(pdf_file_path)

                # 🔹 Extract slide text
                file_info["text_preview"] = extract_pdf_text(pdf_file_path)

                # ✅ Cleanup
                os.remove(temp_file_path)
                os.remove(pdf_file_path)

            except Exception as e:
                file_info["Error"] = f"Failed to process .ppt file: {str(e)}"
            finally:
                if os.path.exists(temp_file_path):
                   os.remove(temp_file_path)  # ✅ Clean up temp file

    return file_info


def count_pdf_pages(pdf_path):
    """Counts pages in the converted PDF (each page = one slide)"""
    from PyPDF2 import PdfReader

    try:
        with open(pdf_path, "rb") as f:
            reader = PdfReader(f)
            return len(reader.pages)
    except:
        return "Unknown"
    
def extract_pdf_text(pdf_path):
    """Extracts text from each slide of the converted PDF"""
    try:
        with open(pdf_path, "rb") as f:
            reader = PdfReader(f)
            text = "\n".join([page.extract_text() for page in reader.pages if page.extract_text()])
            return text[:10000] if text else "No text found"  # Limit to 10000 chars for preview
    except:
        return "Error extracting text"